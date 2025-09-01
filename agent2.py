import os
import requests
import base64
from typing import Dict, Any, List
from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from langchain_community.retrievers import BM25Retriever
from langchain.retrievers import EnsembleRetriever
from smolagents import CodeAgent, OpenAIServerModel, Tool, ToolCallingAgent
from smolagents import PythonInterpreterTool

# Langfuse observability imports
from opentelemetry.sdk.trace import TracerProvider
from openinference.instrumentation.smolagents import SmolagentsInstrumentor
from opentelemetry.exporter.otlp.proto.http.trace_exporter import OTLPSpanExporter
from opentelemetry.sdk.trace.export import SimpleSpanProcessor
from opentelemetry import trace
from langfuse import Langfuse
from smolagents import PythonInterpreterTool


import requests
from markdownify import markdownify
from requests.exceptions import RequestException
from smolagents import tool
import re
import mimetypes
import tempfile
import subprocess
import shutil

import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import pypdf


from concurrent.futures import ThreadPoolExecutor, TimeoutError

from smolagents import Tool
from google import genai
from google.genai import types
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import os
import mimetypes
from typing import Dict, Any

agent_type = "ToolAgent"

client = genai.Client(api_key=os.environ.get("GOOGLE_API_KEY"))


# use GoogleGenerativeAIEmbeddings from langchain_google_genai

def llm_reformat(response: str, question: str, type: str = "number") -> str:
    """
    Use LLM to extract ONLY the number from the response, using the question context to ensure the correct number is chosen.
    """
    if type != "number":
        return response  # Only process if type is number

    format_prompt = f"""Given the following question and response, extract ONLY the number that directly answers the question, following GAIA formatting rules.


Examples:
Question: "How many papers were published in total?"
Response: "The analysis shows 156 papers were published in total."
Answer: 156

Question: "How many stars are there?"
Response: "There are approximately 3.14e+8 stars."
Answer: 3.14e+8

Now, given:
Question: {question}
Response: {response}

Extract ONLY the number that answers the question:
Answer:"""

    try:
        formatting_response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[format_prompt]
        )
        answer = formatting_response.text.strip()
        # Remove any accidental prefix
        if "Answer:" in answer:
            answer = answer.split("Answer:")[-1].strip()
        return answer
    except Exception as e:
        print(f"LLM reformatting failed: {e}")
        return response
    
class FinalAnswerTool(Tool):
    name = "final_answer"
    description = "Provides a final answer to the given problem with automatic formatting based on question keywords."
    inputs = {
        "answer": {
            "type": "any", 
            "description": "The final answer to the problem"
        },
        "original_question": {
            "type": "string", 
            "description": "The original question to determine answer format", 
            "nullable": True  # Optional parameter
        }
    }
    output_type = "any"
    
    def __init__(self):
        super().__init__()
        self.keywords = {
            "how many": "number",
            "how much": "number", 
            "count": "number",
            "total": "number",
            "sum": "number",
            "average": "number",
            "median": "number",
            "percentage": "number",
            "number" : "string",
        }
    
    def _detect_keyword(self, question: str) -> str:
        """Detect which keyword pattern matches the question."""
        question_lower = question.lower().strip()
        
        for keyword, format_type in self.keywords.items():
            if keyword in question_lower:
                return format_type
        
        return "string"  # Default format
    
    def _format_answer(self, answer: Any, format_type: str, original_question) -> Any:
        """Format the answer based on the detected keyword type."""
        print(f"Detected format type: {format_type} for question: {original_question}")
        if format_type == "number":
            # Try to extract/convert to number
            if isinstance(answer, (int, float)):
                return answer
            elif isinstance(answer, str):
                # Extract first number from string
                answer2 = llm_reformat(answer, original_question, type="number")
                return answer2 if answer2 else answer
        return answer
    
    def forward(self, answer: Any, original_question: str = "") -> Any:
        """Process the answer with automatic formatting based on question keywords."""
        
        if not original_question:
            return answer
        
        # Detect the expected format from the question
        format_type = self._detect_keyword(original_question)
        
        # Format the answer accordingly
        formatted_answer = self._format_answer(answer, format_type, original_question)
        
        return formatted_answer

@tool
def get_youtube_transcript(youtube_url: str) -> str:
    """
    Fetches the transcript of a YouTube video given its URL, if available.

    Args:
        youtube_url: The URL of the YouTube video to fetch the transcript for.

    Returns:
        The transcript text if available, or an error message if not.
    """
    try:
        from youtube_transcript_api._api import YouTubeTranscriptApi
        from youtube_transcript_api._errors import TranscriptsDisabled, NoTranscriptFound
        import re
        # Extract video ID from URL
        match = re.search(r"(?:v=|youtu.be/)([\w-]{11})", youtube_url)
        if not match:
            return "Could not extract video ID from the provided URL."
        video_id = match.group(1)
        # Fetch transcript
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        # Combine transcript text
        transcript_text = " ".join([entry['text'] for entry in transcript])
        return transcript_text
    except ImportError:
        return "You must install the 'youtube-transcript-api' package to use this tool. Try: pip install youtube-transcript-api."
    except TranscriptsDisabled:
        return "Transcripts are disabled for this video."
    except NoTranscriptFound:
        return "No transcript found for this video."
    except Exception as e:
        return f"Error fetching transcript: {str(e)}"

class WebSearchTool(Tool):
    name = "web_search"
    description = """Performs a duckduckgo web search based on your query (think a Google search) then returns the top search results."""
    inputs = {"query": {"type": "string", "description": "The search query to perform."}}
    output_type = "string"

    def __init__(self, max_results=10, **kwargs):
        super().__init__()
        self.max_results = max_results
        try:
            from duckduckgo_search import DDGS
        except ImportError as e:
            raise ImportError(
                "You must install package `duckduckgo_search` to run this tool: for instance run `pip install duckduckgo-search`."
            ) from e
        self.ddgs = DDGS(**kwargs)

    def _perform_search(self, query: str):
        """Internal method to perform the actual search."""
        return self.ddgs.text(query, max_results=self.max_results)

    def forward(self, query: str) -> str:
        results = []
        
        # First attempt with timeout
        with ThreadPoolExecutor(max_workers=1) as executor:
            try:
                future = executor.submit(self._perform_search, query)
                results = future.result(timeout=30)  # 30 second timeout
            except TimeoutError:
                print("First search attempt timed out after 30 seconds, retrying...")
                results = []
        
        # Retry if no results or timeout occurred
        if len(results) == 0:
            print("Retrying search...")
            with ThreadPoolExecutor(max_workers=1) as executor:
                try:
                    future = executor.submit(self._perform_search, query)
                    results = future.result(timeout=30)  # 30 second timeout for retry
                except TimeoutError:
                    raise Exception("Search timed out after 30 seconds on both attempts. Try a different query.")
        
        # Final check for results
        if len(results) == 0:
            raise Exception("No results found after two attempts! Try a less restrictive/shorter query.")
        
        postprocessed_results = [f"[{result['title']}]({result['href']})\n{result['body']}" for result in results]
        return "## Search Results\n\n" + "\n\n".join(postprocessed_results)

@tool
def visit_webpage(url: str) -> str:
    """Visits a webpage at the given URL and returns its content as a markdown string.

    Args:
        url: The URL of the webpage to visit.

    Returns:
        The content of the webpage converted to Markdown, or an error message if the request fails.
    """
    try:
        # Send a GET request to the URL
        response = requests.get(url)
        response.raise_for_status()  # Raise an exception for bad status codes

        # Parse the content as HTML with BeautifulSoup
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(response.content, 'html.parser')
        # Extract text and convert to Markdown
        content = soup.get_text(separator="\n", strip=True)
        markdown_content = markdownify(content)
        # Clean up the markdown content
        markdown_content = re.sub(r'\n+', '\n', markdown_content)  # Remove excessive newlines
        markdown_content = re.sub(r'\s+', ' ', markdown_content)  # Remove excessive spaces
        markdown_content = markdown_content.strip()  # Strip leading/trailing whitespace
        return markdown_content

    except RequestException as e:
        return f"Error fetching the webpage: {str(e)}"
    except Exception as e:
        return f"An unexpected error occurred: {str(e)}"
    

class UnifiedMultimodalTool(Tool):
    name = "multimodal_processor"
    description = """
    Unified tool for processing audio, video, and image files.
    Supports transcription, analysis, content extraction, captioning, and cross-modal tasks.
    Handles common formats: mp3, wav, m4a, mp4, avi, mov, jpg, png, gif, etc.
    """
    
    # Required: Define inputs attribute
    inputs = {
        "file_path": {
            "type": "string", 
            "description": "Path to the media file (audio, video, or image) to process"
        },
        "task": {
            "type": "string",
            "description": "Processing task: analyze, transcribe, extract, caption, summarize, or search",
            "nullable": True  # Optional parameter
        },
        "modality": {
            "type": "string", 
            "description": "Force specific modality: auto, audio, video, or image",
            "nullable": True  # Optional parameter
        },
        "additional_context": {
            "type": "string",
            "description": "Extra instructions for processing",
            "nullable": True  # Optional parameter
        }
    }
    
    # Required: Define output type
    output_type = "string"
    
    def __init__(self, api_key: str):
        super().__init__()
        self.client = genai.Client(api_key=api_key)
        # Initialize mimetypes for better detection
        mimetypes.init()
        
    def forward(self, 
                file_path: str, 
                task: str = "analyze", 
                modality: str = "auto",
                additional_context: str = "") -> str:
        """
        Process multimedia files with unified interface
        
        Args:
            file_path: Path to media file
            task: Processing task (analyze, transcribe, extract, caption, summarize, search)
            modality: Force specific modality (auto, audio, video, image)
            additional_context: Extra instructions for processing
        """
        try:
            # Auto-detect modality if not specified
            if modality == "auto":
                modality = self._detect_modality(file_path)
            
            # Check file size for upload method selection
            file_size = os.path.getsize(file_path)
            use_files_api = file_size > 20 * 1024 * 1024  # 20MB threshold
            
            # Generate appropriate prompt based on task and modality
            prompt = self._generate_prompt(task, modality, additional_context)
            
            # Process file
            if use_files_api:
                return self._process_with_files_api(file_path, prompt)
            else:
                return self._process_inline(file_path, prompt)
                
        except Exception as e:
            return f"Error processing {modality} file: {str(e)}"
    
    def _detect_modality(self, file_path: str) -> str:
        """Auto-detect file modality using mimetypes.guess_type()"""
        # Use mimetypes.guess_type for reliable MIME type detection
        mime_type, encoding = mimetypes.guess_type(file_path)
        
        if mime_type:
            if mime_type.startswith('audio/'):
                return 'audio'
            elif mime_type.startswith('video/'):
                return 'video'
            elif mime_type.startswith('image/'):
                return 'image'
        
        # Fallback to extension-based detection if MIME type is unknown
        ext = file_path.lower().split('.')[-1]
        audio_exts = {'mp3', 'wav', 'm4a', 'ogg', 'flac', 'aac', 'wma'}
        video_exts = {'mp4', 'avi', 'mov', 'mkv', 'webm', 'flv', 'wmv', '3gp'}
        image_exts = {'jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp', 'tiff', 'svg'}
        
        if ext in audio_exts:
            return 'audio'
        elif ext in video_exts:
            return 'video'
        elif ext in image_exts:
            return 'image'
        else:
            return 'unknown'
    
    def _get_mime_type(self, file_path: str) -> str:
        """Get MIME type using mimetypes.guess_type()"""
        mime_type, encoding = mimetypes.guess_type(file_path)
        
        if mime_type:
            return mime_type
        
        # Fallback mappings for common formats not detected by mimetypes
        ext = file_path.lower().split('.')[-1]
        fallback_mappings = {
            'm4a': 'audio/mp4',
            'mkv': 'video/x-matroska',
            'webm': 'video/webm',
            'flac': 'audio/flac',
            'ogg': 'audio/ogg',
            'webp': 'image/webp'
        }
        
        return fallback_mappings.get(ext, 'application/octet-stream')
    
    def _generate_prompt(self, task: str, modality: str, context: str) -> str:
        """Generate contextual prompts based on task and modality"""
        
        base_prompts = {
            'analyze': {
                'audio': 'Analyze this audio file. Describe the content, identify sounds, speech, music, and any notable audio characteristics.',
                'video': 'Analyze this video comprehensively. Describe visual content, actions, audio elements, and their relationship.',
                'image': 'Analyze this image in detail. Describe objects, scenes, text, colors, composition, and any notable features.'
            },
            'transcribe': {
                'audio': 'Transcribe all speech in this audio file accurately. Include speaker changes if multiple speakers.',
                'video': 'Transcribe all speech and dialogue in this video. Note visual context when relevant.',
                'image': 'Extract and transcribe any text visible in this image using OCR.'
            },
            'extract': {
                'audio': 'Extract key information, topics, or specific content from this audio.',
                'video': 'Extract key visual and audio information from this video.',
                'image': 'Extract all text, objects, and important visual elements from this image.'
            },
            'caption': {
                'audio': 'Generate descriptive captions for this audio content.',
                'video': 'Generate detailed captions describing both visual and audio elements of this video.',
                'image': 'Generate a comprehensive caption describing this image.'
            },
            'summarize': {
                'audio': 'Provide a concise summary of the main points in this audio.',
                'video': 'Summarize the key visual and audio content of this video.',
                'image': 'Summarize the main elements and content of this image.'
            },
            'search': {
                'audio': 'Make this audio content searchable by extracting keywords, topics, and semantic information.',
                'video': 'Extract searchable content from both visual and audio elements of this video.',
                'image': 'Extract searchable keywords and descriptions from this image.'
            }
        }
        
        prompt = base_prompts.get(task, base_prompts['analyze']).get(modality, 'Analyze this media file.')
        
        if context:
            prompt += f" Additional context: {context}"
            
        return prompt
    
    def _process_with_files_api(self, file_path: str, prompt: str) -> str:
        """Process large files using Files API"""
        uploaded_file = self.client.files.upload(file=file_path)
        
        response = self.client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[prompt, uploaded_file]
        )
        
        return response.text
    
    def _process_inline(self, file_path: str, prompt: str) -> str:
        """Process smaller files inline"""
        with open(file_path, "rb") as f:
            file_data = f.read()
        
        mime_type = self._get_mime_type(file_path)
        
        # Correct way to create parts in the new SDK
        contents = [
            prompt,  # Text can be passed directly
            types.Part.from_bytes(
                data=file_data,
                mime_type=mime_type
            )
        ]
        
        response = self.client.models.generate_content(
            model="gemini-2.5-flash",
            contents=contents
        )
        
        return response.text

    def get_file_info(self, file_path: str) -> Dict[str, str]:
        """Get detailed file information including MIME type and modality"""
        mime_type, encoding = mimetypes.guess_type(file_path)
        modality = self._detect_modality(file_path)
        file_size = os.path.getsize(file_path)
        
        return {
            'file_path': file_path,
            'mime_type': mime_type or 'unknown',
            'encoding': encoding or 'none',
            'modality': modality,
            'file_size': f"{file_size:,} bytes",
            'processing_method': 'Files API' if file_size > 20 * 1024 * 1024 else 'Inline'
        }

class ChromaBM25HybridRetrieverTool(Tool):
    """
    Chroma + BM25 hybrid retriever tool.
    Uses Chroma for dense vector similarity and BM25 for lexical matching.
    """
    name = "chroma_bm25_hybrid_retriever"
    description = "Retrieves relevant document sections using Chroma (dense) and BM25 (lexical)."
    inputs = {
        "query": {
            "type": "string",
            "description": "The search query to find relevant document sections.",
        }
    }
    output_type = "string"

    def __init__(self, chroma: Chroma = None, bm25 = None, top_k: int = 5, alpha: float = 0.5, **kwargs):
        """alpha: weight for dense (Chroma) scores. Final score = alpha * dense + (1-alpha) * lexical"""
        super().__init__(**kwargs)
        self.chroma = chroma
        self.bm25 = bm25
        self.top_k = top_k
        self.alpha = float(alpha)
        # Build an EnsembleRetriever from the provided retrievers. Assume both retrievers are available.
        chroma_retriever = None
        bm25_retriever = None
        if self.chroma:
            try:
                # Chroma provides an as_retriever helper
                chroma_retriever = self.chroma.as_retriever(search_kwargs={"k": self.top_k})
            except Exception:
                # fallback to using similarity methods via a small adapter (not needed per user instruction)
                chroma_retriever = None

        if self.bm25:
            bm25_retriever = self.bm25

        # Create EnsembleRetriever with weights: dense (chroma) weight = alpha, lexical (bm25) weight = 1-alpha
        retrievers = [r for r in (chroma_retriever, bm25_retriever) if r is not None]
        weights = []
        if chroma_retriever is not None:
            weights.append(self.alpha)
        if bm25_retriever is not None:
            weights.append(1.0 - self.alpha)

        # Construct the ensemble retriever
        self.ensemble = EnsembleRetriever(retrievers=retrievers, weights=weights) if retrievers else None

    def _doc_uid(self, d: Document):
        if d.metadata and isinstance(d.metadata, dict):
            return (d.metadata.get('source'), d.page_content[:200])
        return ('', d.page_content[:200])

    def forward(self, query: str) -> str:
        if not self.ensemble:
            return "No documents loaded for retrieval."

        assert isinstance(query, str), "Your search query must be a string"

        # Use EnsembleRetriever to get relevant documents. Slice to top_k to keep behaviour consistent.
        docs = self.ensemble.get_relevant_documents(query)
        if not docs:
            return "No relevant documents found."

        results = docs[: self.top_k]

        return "\nRetrieved documents:\n" + "".join([
            f"\n\n===== Document {str(i)} =====\n" + doc.page_content
            for i, doc in enumerate(results)
        ])

class GAIAAgent:
    """
    GAIA agent using smolagents with Gemini 2.0 Flash and Langfuse observability
    """

    def __init__(self, user_id: str = None, session_id: str = None):
        """Initialize the agent with Gemini 2.0 Flash, tools, and Langfuse observability"""

        # Get API keys
        gemini_api_key = os.environ.get("GOOGLE_API_KEY")
        if not gemini_api_key:
            raise ValueError("GOOGLE_API_KEY environment variable not found")

        # Initialize Langfuse observability
        self._setup_langfuse_observability()

        # Initialize Gemini 2.5 model
        self.model = OpenAIServerModel(
            model_id="gemini-2.5-flash",
            api_base="https://generativelanguage.googleapis.com/v1beta/openai/",
            api_key=gemini_api_key,
            temperature=0.0, 
            top_p=1.0,
        )

        # Store user and session IDs for tracking
        self.user_id = user_id or "gaia-user"
        self.session_id = session_id or "gaia-session"


        # Add a hard coded list of key words in the question and expected format of the answer
        # This is used to help the agent understand the question and expected answer format

        # GAIA system prompt from the leaderboard
        self.system_prompt = f"""You are a general AI assistant. I will ask you a question. Report your thoughts. 
                IMPORTANT:
                - In the last step of your reasoning, if you think your reasoning is not able to answer the question, answer the question directy with your internal reasoning, without using the BM25 retriever tool or the visit_webpage tool.
                - Finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]. YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings. If you are asked for a number, don't use comma to write your number neither use units such as $ or percent sign unless specified otherwise. If you are asked for a string, don't use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise. If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.
                """

    # Initialize retriever tool (will be updated when documents are loaded)
        self.retriever_tool = ChromaBM25HybridRetrieverTool()

        # Create the agent
        self.agent = None
        self._create_agent()

        # Initialize Langfuse client
        self.langfuse = Langfuse()

        from langfuse import get_client
        self.langfuse = get_client()  # ✅ Use get_client() for v3

    def _setup_langfuse_observability(self):
        """Set up Langfuse observability with OpenTelemetry"""
        # Get Langfuse keys from environment variables
        langfuse_public_key = os.environ.get("LANGFUSE_PUBLIC_KEY")
        langfuse_secret_key = os.environ.get("LANGFUSE_SECRET_KEY")
        
        if not langfuse_public_key or not langfuse_secret_key:
            print("Warning: LANGFUSE_PUBLIC_KEY or LANGFUSE_SECRET_KEY not found. Observability will be limited.")
            return

        # Set up Langfuse environment variables
        os.environ["LANGFUSE_HOST"] = os.environ.get("LANGFUSE_HOST", "https://cloud.langfuse.com")
        
        langfuse_auth = base64.b64encode(
            f"{langfuse_public_key}:{langfuse_secret_key}".encode()
        ).decode()
        
        os.environ["OTEL_EXPORTER_OTLP_ENDPOINT"] = os.environ.get("LANGFUSE_HOST") + "/api/public/otel"
        os.environ["OTEL_EXPORTER_OTLP_HEADERS"] = f"Authorization=Basic {langfuse_auth}"

        # Create a TracerProvider for OpenTelemetry
        trace_provider = TracerProvider()
        
        # Add a SimpleSpanProcessor with the OTLPSpanExporter to send traces
        trace_provider.add_span_processor(SimpleSpanProcessor(OTLPSpanExporter()))
        
        # Set the global default tracer provider
        trace.set_tracer_provider(trace_provider)
        self.tracer = trace.get_tracer(__name__)
        
        # Instrument smolagents with the configured provider
        SmolagentsInstrumentor().instrument(tracer_provider=trace_provider)

    def _create_agent(self):
        """Create the CodeAgent with tools"""
        if agent_type == "ToolAgent":
            # Create a ToolCallingAgent if using tool calling

            coder_agent = CodeAgent(
                name = "coder_agent",
                tools=[
                    PythonInterpreterTool(),
                    FinalAnswerTool()],
                model=self.model,
                description="You are a highly skilled coding specialist agent. Your mission is to solve programming, mathematical, and analytical problems end-to-end.",
                max_steps=5,
                additional_authorized_imports=[
                    "math",             # basic calculations
                    "statistics",       # common numeric helpers
                    "itertools",        # safe functional helpers
                    "datetime",         # date handling
                    "random",           # simple randomness (no os access)
                    "re",               # regular expressions
                    "json",             # serialisation / parsing
                ]
            )
            retriever_agent = ToolCallingAgent(
                name = "retriever_agent",
                tools=[self.retriever_tool],
                model=self.model,
                description="You are a specialized document retrieval agent. Your mission is to retrieve relevant parts of documents based on user queries.",
                max_steps=5
            )
            self.agent = ToolCallingAgent(
                tools=[
                    WebSearchTool(),
                    visit_webpage,  # Custom tool for visiting webpages
                    FinalAnswerTool(),
                    get_youtube_transcript, 
                    UnifiedMultimodalTool(api_key=os.environ.get("GOOGLE_API_KEY"))],
                model=self.model,
                description=self.system_prompt,
                max_steps=6, 
                managed_agents = [coder_agent, retriever_agent])
        else : 
            self.agent = CodeAgent(
                tools= [
                    WebSearchTool(),
                    visit_webpage,  # Custom tool for visiting webpages
                    self.retriever_tool,
                    get_youtube_transcript,
                    PythonInterpreterTool(),
                    FinalAnswerTool(), 
                    UnifiedMultimodalTool(api_key=os.environ.get("GOOGLE_API_KEY"))],
                model=self.model,
                description=self.system_prompt, 
                max_steps=5, 
                additional_authorized_imports = [
        "math",             # basic calculations
        "statistics",       # common numeric helpers
        "itertools",        # safe functional helpers
        "datetime",         # date handling
        "random",           # simple randomness (no os access)
        "re",               # regular expressions
        "json",             # serialisation / parsing
    ]       )


    def load_documents_from_file(self, file_path: str):
        """Load and process text documents for BM25+Chroma retrieval, or return raw content for media files.

        Handles common office formats (xlsx, xls, csv, docx, doc, pptx, ppt, pdf) and falls back to LibreOffice
        conversion for legacy binaries. Extracted text is chunked, embedded with Google embeddings, stored in
        Chroma (if available) and a BM25 retriever is created for lexical matches. The hybrid retriever is
        attached to the agent (`ChromaBM25HybridRetrieverTool`) and the agent is recreated so tools are updated.
        """
        try:
            # Lazy optional imports

            mime_type, _ = mimetypes.guess_type(file_path)
            if mime_type is None:
                mime_type = ""
            print(f"Detected MIME type: {mime_type} for file {file_path}")

            # Binary media -> return bytes
            if mime_type.startswith("image") or mime_type.startswith("video") or mime_type.startswith("audio"):
                with open(file_path, "rb") as f:
                    return f.read()

            ext = os.path.splitext(file_path)[1].lower()
            text = None

            # Spreadsheets
            if ext in {".csv", ".xls", ".xlsx"}:
                if pd is None:
                    raise RuntimeError("pandas is required to parse spreadsheets. Install with `pip install pandas openpyxl xlrd`.")
                if ext == ".csv":
                    df = pd.read_csv(file_path)
                    text = df.to_csv(index=False)
                else:
                    sheets = pd.read_excel(file_path, sheet_name=None)
                    parts = []
                    for sheet_name, df in sheets.items():
                        parts.append(f"--- Sheet: {sheet_name} ---")
                        parts.append(df.to_csv(index=False))
                    text = "\n".join(parts)
                    print(text)

            # Word documents
            elif ext in {".docx"}:
                doc = DocxDocument(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
                text = "\n\n".join(paragraphs)

            # PowerPoint
            elif ext in {".pptx"}:
                pres = Presentation(file_path)
                slides_text = []
                for i, slide in enumerate(pres.slides):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text_parts.append(shape.text)
                    slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text_parts))
                text = "\n\n".join(slides_text)

            # Legacy Office binary (.doc, .ppt) or when above parsers unavailable: convert to PDF using LibreOffice
            elif ext in {".doc", ".ppt"}:
                soffice = shutil.which("soffice") or shutil.which("libreoffice")
                if not soffice:
                    raise RuntimeError("LibreOffice (`soffice`) is required to convert legacy Office files. Install LibreOffice or convert to PDF manually.")
                outdir = tempfile.mkdtemp()
                subprocess.run([soffice, "--headless", "--convert-to", "pdf", file_path, "--outdir", outdir],
                               check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                base = os.path.basename(file_path)
                pdfname = os.path.splitext(base)[0] + ".pdf"
                candidate = os.path.join(outdir, pdfname)
                if not os.path.exists(candidate):
                    raise RuntimeError("LibreOffice conversion produced no PDF. Check the input file and LibreOffice installation.")
                if pypdf is None:
                    raise RuntimeError("pypdf is required to extract text from converted PDF. Install with `pip install pypdf`.")
                reader = pypdf.PdfReader(candidate)
                pages = [p.extract_text() or "" for p in reader.pages]
                text = "\n".join(pages)

            # PDF
            elif ext == ".pdf":
                reader = pypdf.PdfReader(file_path)
                pages = [p.extract_text() or "" for p in reader.pages]
                text = "\n".join(pages)

            # Plain text fallback
            else:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()

            if text is None:
                raise RuntimeError("No text could be extracted from the provided file.")

            # Chunking: semantic-first strategy
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                separators=["\n\n", "\n", ".", " ", ""]
            )
            chunks = text_splitter.split_text(text)
            docs = [Document(page_content=chunk, metadata={"source": file_path}) for chunk in chunks]

            # Embeddings (GoogleGenerativeAIEmbeddings)
            embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")

            # Create Chroma vectorstore from documents
            try:
                chroma = Chroma.from_documents(documents=docs, embedding=embeddings)
            except Exception as e:
                print(f"Failed to create Chroma vectorstore: {e}")
                chroma = None

            # Create BM25 retriever
            try:
                bm25 = BM25Retriever(documents=docs)
            except Exception as e:
                print(f"Failed to create BM25 retriever: {e}")
                bm25 = None

            # Attach the hybrid retriever tool and recreate agent so tools list is updated
            self.retriever_tool = ChromaBM25HybridRetrieverTool(chroma=chroma, bm25=bm25, top_k=5, alpha=0.5)
            self._create_agent()

            print(f"Loaded {len(docs)} document chunks from {file_path} into local Chroma store and BM25 retriever")
            return True

        except Exception as e:
            print(f"Error loading documents from {file_path}: {e}")
            return False


    def download_gaia_file(self, task_id: str, api_url: str = "https://agents-course-unit4-scoring.hf.space") -> str:
        """Download file associated with GAIA task_id and return its path"""
        try:
            response = requests.get(f"{api_url}/files/{task_id}", timeout=30)
            response.raise_for_status()

            # Try to get filename from headers
            print(f"Response headers: {response.headers}")
            content_disp = response.headers.get("content-disposition", "")
            match = re.search(r'filename="(.+)"', content_disp)
            if match:
                filename = match.group(1)
            else:
                # Error
                raise ValueError("Filename not found in response headers")

            # Save the file
            with open(filename, 'wb') as f:
                f.write(response.content)

            print(f"Downloaded file saved as {filename}")
            return os.path.abspath(filename)  # Or just return `filename` for relative path

        except Exception as e:
            print(f"Failed to download file for task {task_id}: {e}")
            return None

    def solve_gaia_question(self, question_data: Dict[str, Any], tags: List[str] = None) -> str:
        """
        Solve a GAIA question with full Langfuse observability
        """
        question = question_data.get("Question", "")
        task_id = question_data.get("task_id", "")
        
        # Prepare tags for observability
        trace_tags = ["gaia-agent", "question-solving"]
        if tags:
            trace_tags.extend(tags)
        if task_id:
            trace_tags.append(f"task-{task_id}")

        # Use SDK v3 context manager approach
        with self.langfuse.start_as_current_span(
            name="GAIA-Question-Solving",
            input={"question": question, "task_id": task_id},
            metadata={
                "model": self.model.model_id,
                "question_length": len(question),
                "has_file": bool(task_id)
            }
        ) as span:
            try:
                # Set trace attributes using v3 syntax
                span.update_trace(
                    user_id=self.user_id,
                    session_id=self.session_id,
                    tags=trace_tags
                )

                # Download and load file if task_id provided
                file_loaded = False
                if task_id:
                    file_path = self.download_gaia_file(task_id)
                    if file_path:
                        file_loaded = self.load_documents_from_file(file_path)
                        print(f"Loaded file for task {task_id}")

                # Prepare the prompt
                prompt = f"""
    Question: {question}
    {f'Task ID: {task_id}' if task_id else ''}
    {f'File loaded: Yes as {file_path}' if file_loaded else 'File loaded: No'}

                """

                print("=== AGENT REASONING ===")
                result = self.agent.run(prompt)
                print("=== END REASONING ===")

                # Update span with result using v3 syntax
                span.update(output={"answer": str(result)})

                return str(result)

            except Exception as e:
                error_msg = f"Error processing question: {str(e)}"
                print(error_msg)
                
                # Log error using v3 syntax
                span.update(
                    output={"error": error_msg},
                    level="ERROR"
                )
                
                return error_msg


    def evaluate_answer(self, question: str, answer: str, expected_answer: str = None) -> Dict[str, Any]:
        """
        Evaluate the agent's answer using LLM-as-a-Judge and optionally compare with expected answer
        """
        evaluation_prompt = f"""
Please evaluate the following answer to a question on a scale of 1-5:

Question: {question}
Answer: {answer}
{f'Expected Answer: {expected_answer}' if expected_answer else ''}

Rate the answer on:
1. Accuracy (1-5)
2. Completeness (1-5) 
3. Clarity (1-5)

Provide your rating as JSON: {{"accuracy": X, "completeness": Y, "clarity": Z, "overall": W, "reasoning": "explanation"}}
        """

        try:
            # Use the same model to evaluate
            evaluation_result = self.agent.run(evaluation_prompt)
            
            # Try to parse JSON response
            import json
            scores = json.loads(evaluation_result)
            return scores
        except json.JSONDecodeError:
            # If JSON parsing fails, return a default structure
            print("Failed to parse evaluation result as JSON. Returning default scores.")
            return {
                "accuracy": 0,
                "completeness": 0,
                "clarity": 0,
                "overall": 0,
                "reasoning": "Could not parse evaluation result"
            }


    def add_user_feedback(self, trace_id: str, feedback_score: int, comment: str = None):
        """
        Add user feedback to a specific trace
        
        Args:
            trace_id: The trace ID to add feedback to
            feedback_score: Score from 0-5 (0=very bad, 5=excellent)
            comment: Optional comment from user
        """
        try:
            self.langfuse.score(
                trace_id=trace_id,
                name="user-feedback",
                value=feedback_score,
                comment=comment
            )
            self.langfuse.flush()
            print(f"User feedback added: {feedback_score}/5")
        except Exception as e:
            print(f"Error adding user feedback: {e}")


# Example usage with observability
if __name__ == "__main__":
    # Set up environment variables (you need to set these)
    # os.environ["GOOGLE_API_KEY"] = "your-gemini-api-key"
    # os.environ["LANGFUSE_PUBLIC_KEY"] = "pk-lf-..."
    # os.environ["LANGFUSE_SECRET_KEY"] = "sk-lf-..."
    
    # Test the agent with observability
    agent = GAIAAgent(
        user_id="test-user-123",
        session_id="test-session-456"
    )

    #Example question
    question_data = {
        "Question": "The attached Excel file contains the sales of menu items for a local fast-food chain. What were the total sales that the chain made from food (not including drinks)? Express your answer in USD with two decimal places.",
        "task_id": "7bd855d8-463d-4ed5-93ca-5fe35145f733"
    }

    file = agent.download_gaia_file(question_data["task_id"], api_url="https://agents-course-unit4-scoring.hf.space"
                                    )
    answer = agent.solve_gaia_question(question_data)
    print(f"Answer: {answer}")

